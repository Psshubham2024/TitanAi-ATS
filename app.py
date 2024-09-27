import streamlit as st
from docx import Document
import PyPDF2
import pptx
import json
from io import BytesIO
import re
import requests

# Define the API endpoint and access token
API_URL = "https://api.psnext.info/api/chat"
PSCHATACCESSTOKEN = "your_access_token_here"

# Function to extract text from a Word document
def extract_text_from_word(docx_file):
    doc = Document(docx_file)
    return '\n'.join([para.text for para in doc.paragraphs])

# Function to extract text from a PDF file
def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    pdf_text = ""
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        pdf_text += page.extract_text()
    return pdf_text

# Function to extract text from a PPT file
def extract_text_from_ppt(ppt_file):
    ppt = pptx.Presentation(ppt_file)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Function to extract text from different file types
def extract_text_from_file(file):
    if file.type == "application/pdf":
        return extract_text_from_pdf(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_word(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_ppt(file)
    else:
        return ""

# Improved Function to call PSNext API for CV matching and rating
def get_cv_match(cv_text, job_description):
    payload = {
        "message": f"Evaluate this CV against the following job description. Provide a detailed rating out of 10, specific feedback, and suggestions for improving the CV to better match the role:\n\nJob Description:\n{job_description}\n\nCV:\n{cv_text}",
        "options": {"model": "gpt35turbo"}
    }
    
    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }
    
    response = requests.post(API_URL, headers=headers, data=json.dumps(payload))
    
    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                result = message.get('content', 'No content returned from the API.')
                match = re.search(r'Rating: (\d+)/10', result)
                if match:
                    rating = int(match.group(1))
                    feedback = result.split('\n', 1)[1] if '\n' in result else ''
                    return rating, feedback
                else:
                    return 0, 'Unable to extract rating from the response.'
        return 0, 'No assistant message found in the API response.'
    else:
        return 0, f"Error: {response.status_code}, {response.text}"

# Improved Function to generate case study questions
def generate_case_study_questions(job_description, years_of_experience, industry, difficulty_level):
    payload = {
        "message": f"Based on the following job description and {years_of_experience} years of experience in the {industry} industry, generate a set of case study questions. Ensure the questions reflect a {difficulty_level} difficulty level and challenge critical thinking. Additionally, provide suggestions for how these questions could assess key competencies relevant to the role:\n\nJob Description:\n{job_description}",
        "options": {"model": "gpt35turbo"}
    }
    
    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }
    
    response = requests.post(API_URL, headers=headers, data=json.dumps(payload))
    
    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                return message.get('content', 'No content returned from the API.')
        return 'No assistant message found in the API response.'
    else:
        return f"Error: {response.status_code}, {response.text}"

# Improved Function to match case study answers
def match_case_study_answers(question, provided_answer):
    payload = {
        "message": f"Evaluate the following case study question and answer. Provide a rating out of 10, detailed feedback on the strengths and weaknesses of the response, and suggestions for improvement:\n\nQuestion:\n{question}\n\nAnswer:\n{provided_answer}",
        "options": {"model": "gpt35turbo"}
    }
    
    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }
    
    response = requests.post(API_URL, headers=headers, data=json.dumps(payload))
    
    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                result = message.get('content', 'No content returned from the API.')
                match = re.search(r'Rating: (\d+)/10', result)
                if match:
                    rating = int(match.group(1))
                    feedback = result.split('\n', 1)[1] if '\n' in result else ''
                    return rating, feedback
                else:
                    return 0, 'Unable to extract rating from the response.'
        return 0, 'No assistant message found in the API response.'
    else:
        return 0, f"Error: {response.status_code}, {response.text}"

# Improved Function to compare two sets of texts (questions and solutions)
def compare_question_solution(question_text, solution_text):
    payload = {
        "message": f"Compare the following question document with the solution document. Provide detailed feedback on the alignment between the two, a rating out of 10, and specific suggestions for improving the solution document based on the question:\n\nQuestion Document:\n{question_text}\n\nSolution Document:\n{solution_text}",
        "options": {"model": "gpt35turbo"}
    }
    
    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }
    
    response = requests.post(API_URL, headers=headers, data=json.dumps(payload))
    
    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                result = message.get('content', 'No content returned from the API.')
                match = re.search(r'Rating: (\d+)/10', result)
                if match:
                    rating = int(match.group(1))
                    feedback = result.split('\n', 1)[1] if '\n' in result else ''
                    return rating, feedback
                else:
                    return 0, 'Unable to extract rating from the response.'
        return 0, 'No assistant message found in the API response.'
    else:
        return 0, f"Error: {response.status_code}, {response.text}"

# Emoji mapping based on rating
def get_rating_emoji(rating):
    if rating >= 9:
        return "🌟 Excellent"
    elif rating >= 7:
        return "👍 Good"
    elif rating >= 5:
        return "👌 Average"
    elif rating >= 3:
        return "🤔 Needs Improvement"
    else:
        return "👎 Poor"

# Main app page
def main_app():
    st.title("CareerQgen - Your AI-Powered Staffing Solution")
    st.subheader("Optimize Your Hiring Process with AI")

    tabs = st.tabs(["CV Matching", "Case Study Generation", "Case Study Evaluation", "Document Comparison"])

    # Tab 1: CV Matching
    with tabs[0]:
        st.header("CV Matching")
        st.write("Upload a CV (PDF or Word) and provide a job description to get a match evaluation.")
        
        uploaded_cv = st.file_uploader("Upload Candidate CV (PDF or Word)", type=["pdf", "docx"], key="cv_upload")
        job_description = st.text_area("Job Description", height=200, key="job_description_cv")
        
        result_output = st.empty()

        if uploaded_cv and job_description:
            if uploaded_cv.type == "application/pdf":
                cv_text = extract_text_from_pdf(uploaded_cv)
            elif uploaded_cv.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                cv_text = extract_text_from_word(uploaded_cv)

            if cv_text.strip():
                if st.button("Get CV Match", key="cv_match_button"):
                    with st.spinner("Processing..."):
                        rating, feedback = get_cv_match(cv_text, job_description)
                        result_output.text_area("Match Feedback", feedback, height=200, key="cv_feedback")
                        st.write(f"Rating: {rating}/10")
                        st.write(f"Evaluation: {get_rating_emoji(rating)}")

    # Tab 2: Case Study Generation
    with tabs[1]:
        st.header("Generate Case Study Questions")
        job_description = st.text_area("Job Description", height=200, key="job_description_case")
        years_of_experience = st.text_input("Years of Experience", key="years_of_experience")
        industry = st.text_input("Industry", key="industry")
        difficulty_level = st.selectbox("Difficulty Level", ["Easy", "Intermediate", "Hard"], key="difficulty_level")
        
        if st.button("Generate Case Study Questions", key="generate_questions"):
            with st.spinner("Generating..."):
                questions = generate_case_study_questions(job_description, years_of_experience, industry, difficulty_level)
                st.text_area("Generated Case Study Questions", questions, height=200, key="generated_questions")

    # Tab 3: Case Study Evaluation
    with tabs[2]:
        st.header("Evaluate Case Study Answers")
        case_study_question = st.text_area("Case Study Question", height=100, key="case_study_question")
        provided_answer = st.text_area("Candidate's Answer", height=200, key="provided_answer")
        
        if st.button("Evaluate Answer", key="evaluate_answer"):
            with st.spinner("Evaluating..."):
                rating, feedback = match_case_study_answers(case_study_question, provided_answer)
                st.write(f"Rating: {rating}/10")
                st.write("Feedback:", feedback)
                st.write(f"Evaluation: {get_rating_emoji(rating)}")

    # Tab 4: Document Comparison
    with tabs[3]:
        st.header("Compare Question and Solution Documents")
        question_doc = st.file_uploader("Upload Question Document (PDF, Word, or PowerPoint)", type=["pdf", "docx", "pptx"], key="question_doc_upload")
        solution_doc = st.file_uploader("Upload Solution Document (PDF, Word, or PowerPoint)", type=["pdf", "docx", "pptx"], key="solution_doc_upload")
        
        if st.button("Compare Documents", key="compare_documents"):
            if question_doc and solution_doc:
                with st.spinner("Comparing..."):
                    question_text = extract_text_from_file(question_doc)
                    solution_text = extract_text_from_file(solution_doc)
                    
                    if question_text and solution_text:
                        rating, feedback = compare_question_solution(question_text, solution_text)
                        st.write(f"Rating: {rating}/10")
                        st.write("Feedback:", feedback)
                        st.write(f"Evaluation: {get_rating_emoji(rating)}")

# Run the app
if __name__ == "__main__":
    main_app()
